#!/usr/bin/env python3
"""Build a 3-panel comparison deck where the third panel is the *live* native
splat shapes (grafted from the native PPTX), not a screenshot of them.

Panels: Original | Canvas splat (raster) | PPTX-native splat (live DrawingML).

The native splat group is copied verbatim and re-anchored at exactly its native
EMU size (scale 1.000) so PowerPoint renders the real shapes faithfully.
"""
from copy import deepcopy

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from lxml import etree

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

NATIVE_DECK = "tmp/chameleon_400px_10m_2000_exact_softedge.pptx"
ORIGINAL_IMG = "input.png"
CANVAS_IMG = "tmp/chameleon_400px_10m_2000_exact_canvas_preview.png"
OUT = "tmp/chameleon_400px_10m_2000_comparison_live.pptx"

# Native slide / splat-group footprint (EMU). Panel 2 uses this exact size -> scale 1.000.
BOX_W = 3609975
BOX_H = 3810000

SLIDE_W = 12192000  # 16:9
SLIDE_H = 6858000

CAPTIONS = ["Original", "Canvas splat  (SSIM 0.89)", "PPTX-native splat — live  (SSIM 0.73)"]


def qn(tag):
    return "{%s}%s" % (P_NS, tag)


def layout():
    """Return (x0, x1, x2, box_y) with 3 native-size boxes evenly spaced."""
    total_boxes = BOX_W * 3
    gap = (SLIDE_W - total_boxes) // 4
    x0 = gap
    x1 = x0 + BOX_W + gap
    x2 = x1 + BOX_W + gap
    box_y = 1_550_000
    return x0, x1, x2, box_y


def add_caption(slide, x, y, w, text):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(360000))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x22, 0x22, 0x22)


def add_title(slide):
    tb = slide.shapes.add_textbox(Emu(0), Emu(180000), Emu(SLIDE_W), Emu(600000))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Chameleon — Original vs Canvas splat vs PPTX-native splat (2000 splats)"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x11, 0x11, 0x11)


def add_footer(slide):
    tb = slide.shapes.add_textbox(Emu(0), Emu(SLIDE_H - 430000), Emu(SLIDE_W), Emu(360000))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "github.com/BramAlkema/SplatThis"
    r.font.size = Pt(11)
    r.font.color.rgb = RGBColor(0x66, 0x66, 0x66)


def graft_native_group(slide, x, y):
    """Copy the 'Splat Group' grpSp from the native deck and re-anchor it."""
    src = Presentation(NATIVE_DECK)
    sp_tree = src.slides[0].shapes._spTree
    grp = None
    for g in sp_tree.findall(qn("grpSp")):
        cnvpr = g.find(qn("nvGrpSpPr") + "/" + qn("cNvPr"))
        if cnvpr is not None and cnvpr.get("name") == "Splat Group":
            grp = g
            break
    if grp is None:
        raise RuntimeError("Splat Group not found in native deck")

    grp = deepcopy(grp)

    # Keep all unique ids distinct from the target slide's shapes.
    for cnvpr in grp.iter(qn("cNvPr")):
        cnvpr.set("id", str(int(cnvpr.get("id")) + 500000))

    # Re-anchor: ext stays native (scale 1.000), only move the origin.
    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    xfrm = grp.find(qn("grpSpPr") + "/" + A + "xfrm")
    off = xfrm.find(A + "off")
    off.set("x", str(int(x)))
    off.set("y", str(int(y)))
    ext = xfrm.find(A + "ext")
    ext.set("cx", str(BOX_W))
    ext.set("cy", str(BOX_H))
    choff = xfrm.find(A + "chOff")
    choff.set("x", "0")
    choff.set("y", "0")
    chext = xfrm.find(A + "chExt")
    chext.set("cx", str(BOX_W))
    chext.set("cy", str(BOX_H))

    slide.shapes._spTree.append(grp)


def main():
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    x0, x1, x2, box_y = layout()
    cap_y = box_y - 380000

    add_title(slide)
    add_footer(slide)
    for x, cap in zip((x0, x1, x2), CAPTIONS):
        add_caption(slide, x, cap_y, BOX_W, cap)

    # Panels 0 & 1: raster images, sized to the exact box.
    slide.shapes.add_picture(ORIGINAL_IMG, Emu(x0), Emu(box_y), Emu(BOX_W), Emu(BOX_H))
    slide.shapes.add_picture(CANVAS_IMG, Emu(x1), Emu(box_y), Emu(BOX_W), Emu(BOX_H))

    # Panel 2: live native splat shapes.
    graft_native_group(slide, x2, box_y)

    prs.save(OUT)
    print("wrote", OUT)


if __name__ == "__main__":
    main()
